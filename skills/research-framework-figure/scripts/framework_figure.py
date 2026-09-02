#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
framework_figure.py — 研究架構圖（概念模型圖）產生器

從一份 JSON 規格產出學術投稿等級的研究架構圖。支援中介、調節、被調節的中介、
序列中介等版式，輸出可再編輯的向量圖。

【設計原則】
1. **零依賴出 SVG**：SVG 由字串直接組出，只用標準庫，任何環境都跑得動。
   PPTX 匯出才需要 python-pptx（可選）。
2. **中英混排字型**：font-family 一律寫成 "Times New Roman, DFKai-SB, serif"——
   西文字型排在前面，中文字缺字時自動 fallback 到標楷體。這正是
   CLAUDE.md 要求的「中文＝標楷體、英數＝Times New Roman」在 SVG 的正確作法。
   ⚠️ 若檢視端沒有標楷體，中文會 fallback 成其他字型；產出後請自行確認，
      不要假裝已套上標楷體（同 management-figure 的誠實紀律）。
3. **可再編輯**：SVG 可用 Illustrator/Inkscape 開啟編修；PPTX 為原生圖形，
   可在 PowerPoint 內直接改字改框（口試前微調用）。

【用法】
    python framework_figure.py spec.json -o fig.svg
    python framework_figure.py spec.json -o fig.pptx --format pptx
    python framework_figure.py --template mediation_dual_controls --demo -o demo.svg
    python framework_figure.py --list-templates

last_verified: 2026-07-26
"""

from __future__ import annotations

import argparse
import json
import sys
from typing import Any

# ── 版面常數（可由 spec 的 style 區塊覆寫）──────────────────────────────
DEFAULTS = {
    "font_family": "Times New Roman, DFKai-SB, BiauKai, serif",
    "font_size_header": 19,     # 方框標題列
    "font_size_item": 17,       # 方框內條列
    "font_size_hypo": 20,       # H1/H2/H3 假說標籤
    "stroke": "#000000",
    "stroke_width": 1.8,
    "fill": "#FFFFFF",
    "text_color": "#000000",
    "line_height": 27,
    "box_pad_x": 14,
    "box_pad_y": 12,
    "header_gap": 8,            # 標題列與內容的分隔線間距
    "margin": 40,
}

TEMPLATES = {
    "mediation_dual_controls": "中介模型＋雙控制變數框（每個依變數各一組控制變數）",
    "mediation": "簡單中介模型 X→M→Y（單一控制變數框，可省略）",
    "moderation": "調節模型 W 調節 X→Y（調節箭頭指向路徑中點）",
    "moderated_mediation": "被調節的中介 X→M→Y，W 調節其中一條路徑",
    "serial_mediation": "序列中介 X→M1→M2→Y",
}


# ── 小工具 ─────────────────────────────────────────────────────────────
def esc(s: Any) -> str:
    """XML 逸出。"""
    return (
        str(s)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def text_w(s: str, size: float) -> float:
    """估算字串寬度：CJK 全形約 1.0em，西文約 0.52em。用於自動決定框寬。"""
    w = 0.0
    for ch in str(s):
        w += size * (1.0 if ord(ch) > 0x2E80 else 0.52)
    return w


class Box:
    """一個帶標題列與條列內容的方框。"""

    def __init__(self, label: str, items: list[str], st: dict, min_w: float = 0):
        self.label = label
        self.items = items or []
        self.st = st
        self.numbered = len(self.items) > 1

        widest = text_w(label, st["font_size_header"])
        for i, it in enumerate(self.items, 1):
            prefix = f"{i}. " if self.numbered else ""
            widest = max(widest, text_w(prefix + it, st["font_size_item"]))
        self.w = max(min_w, widest + st["box_pad_x"] * 2)

        self.header_h = st["font_size_header"] + st["box_pad_y"] * 2
        body = len(self.items) * st["line_height"]
        self.body_h = (body + st["box_pad_y"] * 2) if self.items else 0
        self.h = self.header_h + self.body_h
        self.x = 0.0
        self.y = 0.0

    # 幾何輔助
    @property
    def cx(self) -> float:
        return self.x + self.w / 2

    @property
    def cy(self) -> float:
        return self.y + self.h / 2

    def anchor(self, side: str) -> tuple[float, float]:
        return {
            "top": (self.cx, self.y),
            "bottom": (self.cx, self.y + self.h),
            "left": (self.x, self.cy),
            "right": (self.x + self.w, self.cy),
        }[side]

    def svg(self) -> str:
        st = self.st
        out = [
            f'<rect x="{self.x:.1f}" y="{self.y:.1f}" width="{self.w:.1f}" '
            f'height="{self.h:.1f}" fill="{st["fill"]}" stroke="{st["stroke"]}" '
            f'stroke-width="{st["stroke_width"]}"/>'
        ]
        # 標題列（置中）
        ty = self.y + self.header_h / 2 + st["font_size_header"] * 0.36
        out.append(
            f'<text x="{self.cx:.1f}" y="{ty:.1f}" text-anchor="middle" '
            f'font-family="{esc(st["font_family"])}" font-size="{st["font_size_header"]}" '
            f'font-weight="bold" fill="{st["text_color"]}">{esc(self.label)}</text>'
        )
        if not self.items:
            return "\n".join(out)

        # 標題與內容的分隔線
        sep = self.y + self.header_h
        out.append(
            f'<line x1="{self.x:.1f}" y1="{sep:.1f}" x2="{self.x + self.w:.1f}" '
            f'y2="{sep:.1f}" stroke="{st["stroke"]}" stroke-width="{st["stroke_width"]}"/>'
        )
        # 條列內容：多項編號者靠左（易讀），單項者置中（符合學術架構圖慣例）
        centered = not self.numbered
        tx = self.cx if centered else self.x + st["box_pad_x"]
        anchor = ' text-anchor="middle"' if centered else ""
        for i, it in enumerate(self.items, 1):
            iy = (
                sep
                + st["box_pad_y"]
                + (i - 1) * st["line_height"]
                + st["font_size_item"] * 0.9
            )
            prefix = f"{i}. " if self.numbered else ""
            out.append(
                f'<text x="{tx:.1f}" y="{iy:.1f}"{anchor} font-family="{esc(st["font_family"])}" '
                f'font-size="{st["font_size_item"]}" font-weight="bold" '
                f'fill="{st["text_color"]}">{esc(prefix + it)}</text>'
            )
        return "\n".join(out)


def arrow(p1, p2, st, label="", dashed=False, curve=0.0) -> str:
    """畫一條帶箭頭的線；curve≠0 時畫二次貝茲曲線（控制變數用虛線彎箭頭）。"""
    x1, y1 = p1
    x2, y2 = p2
    dash = ' stroke-dasharray="7,5"' if dashed else ""
    if curve:
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        dx, dy = x2 - x1, y2 - y1
        n = (dx**2 + dy**2) ** 0.5 or 1.0
        cxp, cyp = mx - dy / n * curve, my + dx / n * curve
        d = f"M {x1:.1f},{y1:.1f} Q {cxp:.1f},{cyp:.1f} {x2:.1f},{y2:.1f}"
    else:
        d = f"M {x1:.1f},{y1:.1f} L {x2:.1f},{y2:.1f}"

    # class="conn" 是給 PPTX 匯出用的語意標記：svg2drawingml 會據此挑出
    # 連接線轉成原生 custGeom 形狀。改動請同步 build_pptx()。
    out = [
        f'<path class="conn" d="{d}" fill="none" stroke="{st["stroke"]}" '
        f'stroke-width="{st["stroke_width"]}"{dash} marker-end="url(#ah)"/>'
    ]
    if label:
        lx, ly = (x1 + x2) / 2, (y1 + y2) / 2
        out.append(
            f'<rect x="{lx - 22:.1f}" y="{ly - 26:.1f}" width="44" height="26" '
            f'fill="{st["fill"]}" stroke="none"/>'
        )
        out.append(
            f'<text class="hypo" x="{lx:.1f}" y="{ly - 7:.1f}" text-anchor="middle" '
            f'font-family="{esc(st["font_family"])}" font-size="{st["font_size_hypo"]}" '
            f'font-weight="bold" fill="{st["text_color"]}">{esc(label)}</text>'
        )
    return "\n".join(out)


# ── 版式：中介＋雙控制變數框（使用者附圖的標準版式）────────────────────
def layout_mediation_dual_controls(spec: dict, st: dict) -> tuple[str, float, float]:
    m = Box(spec["m"]["label"], spec["m"].get("items", []), st)
    x = Box(spec["x"]["label"], spec["x"].get("items", []), st, min_w=250)
    y = Box(spec["y"]["label"], spec["y"].get("items", []), st, min_w=330)
    ctrls = [Box(c["label"], c.get("items", []), st) for c in spec.get("controls", [])]

    mg = st["margin"]
    gap_v = 90          # 上排與中排的垂直間距
    gap_h = 110         # X／Y 與中央的水平間距

    top_w = max(m.w, x.w + y.w + gap_h * 2)
    bot_w = sum(c.w for c in ctrls) + (100 if len(ctrls) > 1 else 0)
    W = max(top_w, bot_w) + mg * 2

    # 上排：M 置中
    m.x, m.y = (W - m.w) / 2, mg
    # 中排：X 左、Y 右
    mid_y = m.y + m.h + gap_v
    x.x, x.y = mg + 20, mid_y
    y.x, y.y = W - mg - 20 - y.w, mid_y - 10

    # 下排：控制變數框
    bot_y = max(x.y + x.h, y.y + y.h) + 95
    if ctrls:
        if len(ctrls) == 1:
            ctrls[0].x = (W - ctrls[0].w) / 2
        else:
            ctrls[0].x = mg + 20
            ctrls[1].x = W - mg - 20 - ctrls[1].w
        for c in ctrls:
            c.y = bot_y

    H = (max([c.y + c.h for c in ctrls]) if ctrls else mid_y + max(x.h, y.h)) + mg

    parts = [b.svg() for b in [m, x, y] + ctrls]
    hyp = spec.get("hypotheses", {})

    # H1：X → M（左下往右上）
    parts.append(
        arrow(
            (x.x + x.w * 0.55, x.y),
            (m.x + m.w * 0.06, m.y + m.h),
            st,
            hyp.get("h1", "H1"),
        )
    )
    # H2：X → Y（水平主路徑）
    parts.append(arrow(x.anchor("right"), (y.x, y.cy), st, hyp.get("h2", "H2")))
    # H3：M → Y（右上往右下）
    parts.append(
        arrow(
            (m.x + m.w * 0.94, m.y + m.h),
            (y.x + y.w * 0.45, y.y),
            st,
            hyp.get("h3", "H3"),
        )
    )
    # 控制變數：虛線彎箭頭指向其對應的依變數
    tgt = {"m": m, "y": y}
    for i, c in enumerate(ctrls):
        t = tgt.get(spec["controls"][i].get("target", "y"), y)
        parts.append(
            arrow(
                (c.cx, c.y),
                (t.cx - t.w * 0.18, t.y + t.h) if t is m else (t.cx, t.y + t.h),
                st,
                dashed=True,
                curve=55 if i == 0 else -55,
            )
        )
    return "\n".join(parts), W, H


def layout_generic(spec: dict, st: dict) -> tuple[str, float, float]:
    """中介／調節／被調節的中介／序列中介的通用水平版式。"""
    tpl = spec.get("template", "mediation")
    mg = st["margin"]
    nodes: list[Box] = []

    xb = Box(spec["x"]["label"], spec["x"].get("items", []), st, min_w=230)
    yb = Box(spec["y"]["label"], spec["y"].get("items", []), st, min_w=230)
    mids = spec.get("mediators") or ([spec["m"]] if spec.get("m") else [])
    mbs = [Box(d["label"], d.get("items", []), st, min_w=230) for d in mids]

    gap = 130
    row_h = max([b.h for b in [xb, yb] + mbs] or [80])
    top_y = mg + 150

    if tpl == "serial_mediation" and len(mbs) >= 2:
        seq = [xb] + mbs + [yb]
        cx = mg
        for b in seq:
            b.x, b.y = cx, top_y
            cx += b.w + gap
        W = cx - gap + mg
        nodes = seq
        parts = [b.svg() for b in nodes]
        hyp = spec.get("hypotheses", {})
        for i in range(len(seq) - 1):
            parts.append(
                arrow(
                    seq[i].anchor("right"),
                    seq[i + 1].anchor("left"),
                    st,
                    hyp.get(f"h{i + 1}", ""),
                )
            )
        H = top_y + row_h + mg
        return "\n".join(parts), W, H

    # 中介：X 左、M 上中、Y 右
    mb = mbs[0] if mbs else None
    xb.x, xb.y = mg, top_y
    if mb:
        mb.x = mg + xb.w + gap
        mb.y = mg
        yb.x = mb.x + mb.w + gap
    else:
        yb.x = mg + xb.w + gap * 2
    yb.y = top_y
    W = yb.x + yb.w + mg
    nodes = [b for b in [xb, mb, yb] if b]
    parts = [b.svg() for b in nodes]
    hyp = spec.get("hypotheses", {})

    if mb:
        parts.append(
            arrow((xb.cx, xb.y), (mb.x, mb.y + mb.h), st, hyp.get("h1", "H1"))
        )
        parts.append(
            arrow((mb.x + mb.w, mb.y + mb.h), (yb.cx, yb.y), st, hyp.get("h3", "H3"))
        )
    parts.append(arrow(xb.anchor("right"), yb.anchor("left"), st, hyp.get("h2", "H2")))

    H = top_y + row_h + mg
    # 調節變數：箭頭指向 X→Y 路徑中點（學術慣例，不是指向 Y）
    if spec.get("moderator"):
        wb = Box(
            spec["moderator"]["label"], spec["moderator"].get("items", []), st, min_w=230
        )
        wb.x, wb.y = (xb.cx + yb.cx) / 2 - wb.w / 2, top_y + row_h + 110
        parts.append(wb.svg())
        midp = ((xb.x + xb.w + yb.x) / 2, xb.cy)
        parts.append(arrow((wb.cx, wb.y), midp, st))
        H = wb.y + wb.h + mg
    return "\n".join(parts), W, H


def build_svg(spec: dict) -> str:
    st = dict(DEFAULTS)
    st.update(spec.get("style", {}))
    tpl = spec.get("template", "mediation")

    if tpl == "mediation_dual_controls":
        body, W, H = layout_mediation_dual_controls(spec, st)
    else:
        body, W, H = layout_generic(spec, st)

    title = spec.get("title", "")
    title_svg = ""
    if title:
        title_svg = (
            f'<text x="{W / 2:.1f}" y="{st["margin"] * 0.75:.1f}" text-anchor="middle" '
            f'font-family="{esc(st["font_family"])}" font-size="{st["font_size_header"] + 4}" '
            f'font-weight="bold" fill="{st["text_color"]}">{esc(title)}</text>'
        )
        H += 10

    return f"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W:.0f} {H:.0f}" width="{W:.0f}" height="{H:.0f}">
<defs>
<marker id="ah" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="7" markerHeight="7" orient="auto-start-reverse">
<path d="M 0 0 L 10 5 L 0 10 z" fill="{st['stroke']}"/>
</marker>
</defs>
<rect width="100%" height="100%" fill="#FFFFFF"/>
{title_svg}
{body}
</svg>
"""


# ── PPTX 匯出（可選，需 python-pptx）────────────────────────────────────
def build_pptx(spec: dict, out_path: str) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import qn
        from lxml import etree  # noqa: F401  (python-pptx 依賴，確認可用)
    except ImportError:
        raise SystemExit(
            "錯誤：PPTX 匯出需要 python-pptx。請執行：pip install python-pptx\n"
            "（SVG 輸出不需要任何套件，可改用 --format svg）"
        )

    st = dict(DEFAULTS)
    st.update(spec.get("style", {}))

    # 以 SVG 版面為基準換算：1 SVG 單位 = 9525 EMU * scale
    scale = 9525 * 0.75

    def set_font(run, size, bold=True):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Times New Roman"
        run.font.color.rgb = RGBColor(0, 0, 0)
        # 東亞字型必須另外寫 a:ea，否則中文不會套標楷體
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn("a:ea"), {"typeface": "標楷體"})
        rPr.append(ea)

    prs = Presentation()
    prs.slide_width = Emu(int(13.333 * 914400))
    prs.slide_height = Emu(int(7.5 * 914400))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tpl = spec.get("template", "mediation")
    if tpl == "mediation_dual_controls":
        _, W, H = layout_mediation_dual_controls(spec, st)
    else:
        _, W, H = layout_generic(spec, st)

    # 重建 Box 幾何以放置 PPTX 圖形
    boxes: list[Box] = []
    if tpl == "mediation_dual_controls":
        m = Box(spec["m"]["label"], spec["m"].get("items", []), st)
        xb = Box(spec["x"]["label"], spec["x"].get("items", []), st, min_w=250)
        yb = Box(spec["y"]["label"], spec["y"].get("items", []), st, min_w=330)
        cs = [Box(c["label"], c.get("items", []), st) for c in spec.get("controls", [])]
        mg, gap_v, gap_h = st["margin"], 90, 110
        top_w = max(m.w, xb.w + yb.w + gap_h * 2)
        bot_w = sum(c.w for c in cs) + (100 if len(cs) > 1 else 0)
        Wp = max(top_w, bot_w) + mg * 2
        m.x, m.y = (Wp - m.w) / 2, mg
        mid_y = m.y + m.h + gap_v
        xb.x, xb.y = mg + 20, mid_y
        yb.x, yb.y = Wp - mg - 20 - yb.w, mid_y - 10
        bot_y = max(xb.y + xb.h, yb.y + yb.h) + 95
        if cs:
            if len(cs) == 1:
                cs[0].x = (Wp - cs[0].w) / 2
            else:
                cs[0].x, cs[1].x = mg + 20, Wp - mg - 20 - cs[1].w
            for c in cs:
                c.y = bot_y
        boxes = [m, xb, yb] + cs

    # ── 箭頭與假說標籤：由 SVG 轉成原生 custGeom 形狀 ──────────────
    # python-pptx 畫不出曲線箭頭，因此以 SVG 為單一事實來源，
    # 用 svg2drawingml 做確定性轉換。缺這段會靜默丟失所有箭頭。
    arrows_added = 0
    try:
        import xml.etree.ElementTree as _ET
        import sys as _sys, os as _os
        _sys.path.insert(0, _os.path.dirname(_os.path.abspath(__file__)))
        from svg2drawingml import (
            SVG_NS, parse_path, path_bbox, build_path_shape_xml, scan_svg,
        )

        svg_text = build_svg(spec)
        info = scan_svg(svg_text)
        if info["unsupported"]:
            print(f"注意：SVG 含本轉換器不支援的元素 {info['unsupported']}，"
                  f"該部分不會出現在 PPTX（不靜默略過，特此告知）。")

        root = _ET.fromstring(svg_text)
        spTree = slide.shapes._spTree
        sid = 900
        for el in root.iter():
            tag = el.tag.replace(SVG_NS, "")
            cls = el.get("class", "")
            if tag == "path" and cls == "conn":
                cmds = parse_path(el.get("d", ""))
                if not cmds:
                    continue
                bb = path_bbox(cmds)
                sid += 1
                xml = build_path_shape_xml(
                    sid, f"conn{sid}", cmds, bb,
                    fill="none", stroke=el.get("stroke", "#000000"),
                    width_px=el.get("stroke-width", 1.8),
                    dashed=bool(el.get("stroke-dasharray")),
                    arrow_end=bool(el.get("marker-end")),
                )
                spTree.append(parse_xml(xml))
                arrows_added += 1
            elif tag == "text" and cls == "hypo":
                tx, ty = float(el.get("x", 0)), float(el.get("y", 0))
                tb = slide.shapes.add_textbox(
                    Emu(int((tx - 28) * scale)), Emu(int((ty - 24) * scale)),
                    Emu(int(56 * scale)), Emu(int(28 * scale)))
                tf = tb.text_frame
                tf.word_wrap = False
                tf.paragraphs[0].text = el.text or ""
                if tf.paragraphs[0].runs:
                    set_font(tf.paragraphs[0].runs[0], st["font_size_hypo"] * 0.62)
                arrows_added += 1
    except ImportError as e:
        print(f"警告：svg2drawingml 載入失敗（{e}），本次 PPTX 將不含箭頭。")

    for b in boxes:
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(b.x * scale)),
            Emu(int(b.y * scale)),
            Emu(int(b.w * scale)),
            Emu(int(b.h * scale)),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shp.line.color.rgb = RGBColor(0, 0, 0)
        shp.line.width = Pt(1.5)

        tf = shp.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = b.label
        set_font(tf.paragraphs[0].runs[0], st["font_size_header"] * 0.62)
        numbered = len(b.items) > 1
        for i, it in enumerate(b.items, 1):
            p = tf.add_paragraph()
            p.text = (f"{i}. " if numbered else "") + it
            set_font(p.runs[0], st["font_size_item"] * 0.62)

    prs.save(out_path)


# ── 內建示範規格 ───────────────────────────────────────────────────────
DEMO = {
    "template": "mediation_dual_controls",
    "x": {"label": "自變數(X)", "items": ["〔自變數名稱〕"]},
    "m": {
        "label": "作用機制變數(M)",
        "items": ["〔中介變數1〕", "〔中介變數2〕", "〔中介變數3〕"],
    },
    "y": {"label": "應變數(Y)", "items": ["〔應變數名稱〕"]},
    "hypotheses": {"h1": "H1", "h2": "H2", "h3": "H3"},
    "controls": [
        {
            "label": "當應變數為〔M〕時的控制變數(C)",
            "target": "m",
            "items": ["公司年齡", "公司規模", "董事會規模", "負債比率"],
        },
        {
            "label": "當應變數為〔Y〕時的控制變數(C)",
            "target": "y",
            "items": ["公司年齡", "公司規模", "獨立董事占比", "ROA"],
        },
    ],
}


def main() -> int:
    ap = argparse.ArgumentParser(
        description="研究架構圖（概念模型圖）產生器——輸入 JSON 規格，輸出 SVG／PPTX。",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""範例：
  python framework_figure.py spec.json -o fig.svg
  python framework_figure.py spec.json -o fig.pptx --format pptx
  python framework_figure.py --demo -o demo.svg
  python framework_figure.py --list-templates

規格格式與版式說明見 ../references/framework-figure-spec.md
""",
    )
    ap.add_argument("spec", nargs="?", help="JSON 規格檔路徑")
    ap.add_argument("-o", "--output", help="輸出檔路徑（.svg 或 .pptx）")
    ap.add_argument("--format", choices=["svg", "pptx"], help="輸出格式（預設依副檔名）")
    ap.add_argument("--demo", action="store_true", help="用內建示範規格產圖")
    ap.add_argument("--list-templates", action="store_true", help="列出可用版式")
    a = ap.parse_args()

    if a.list_templates:
        print("可用版式（template 欄位）：")
        for k, v in TEMPLATES.items():
            print(f"  {k:28s} {v}")
        return 0

    if a.demo:
        spec = DEMO
    elif a.spec:
        try:
            with open(a.spec, encoding="utf-8") as f:
                spec = json.load(f)
        except FileNotFoundError:
            print(f"錯誤：找不到規格檔 {a.spec}", file=sys.stderr)
            return 1
        except json.JSONDecodeError as e:
            print(f"錯誤：JSON 格式有誤（{e}）", file=sys.stderr)
            return 1
    else:
        ap.print_help()
        return 1

    out = a.output or "framework.svg"
    fmt = a.format or ("pptx" if out.lower().endswith(".pptx") else "svg")

    for key in ("x", "y"):
        if key not in spec:
            print(f"錯誤：規格缺少必要欄位 '{key}'", file=sys.stderr)
            return 1

    if fmt == "pptx":
        build_pptx(spec, out)
    else:
        with open(out, "w", encoding="utf-8") as f:
            f.write(build_svg(spec))

    print(f"成功：已輸出 {out}（版式 {spec.get('template', 'mediation')}，格式 {fmt}）")
    if fmt == "svg":
        print("提示：SVG 可用瀏覽器預覽、Illustrator/Inkscape 編修，或直接插入 Word。")
        print("　　　中文若未顯示為標楷體，代表檢視端缺該字型，請於有字型的機器重出。")
    return 0


if __name__ == "__main__":
    sys.exit(main())
