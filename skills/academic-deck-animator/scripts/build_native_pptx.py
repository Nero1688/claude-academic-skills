#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_native_pptx.py — slides_content.json → 原生 PPTX(引擎B:native-pptx)

用法:
    python build_native_pptx.py slides_content.json -o deck.pptx          # 含原生動畫
    python build_native_pptx.py slides_content.json -o print.pptx --split # 動畫拆頁(轉PDF用)

依賴:python-pptx(pip install python-pptx)。

動畫實作說明(改動前必讀):
python-pptx 沒有動畫 API,本腳本以 OOXML <p:timing> 時間樹直接注入。
EFFECT_MAP 的 presetID/presetSubtype 對應 PowerPoint 內建進場效果編號,
這些編號是 OOXML 規格值,不要憑記憶新增——先用 PowerPoint 手做該效果、
解壓 pptx 對照 slide XML 確認編號後再加。
trigger="after-previous" 的元素會併入前一個 entry_order 組(以 delay_ms 錯開),
這是刻意的簡化:維持時間樹結構單純、避免產生 PowerPoint 無法解析的巢狀。
"""
import argparse
import json
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):  # Windows 主控台中文輸出
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    sys.exit("[錯誤] 需要 python-pptx:pip install python-pptx")

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# effect 名稱 → (presetID, presetSubtype, 實作函式名)。編號為 OOXML 規格值。
EFFECT_MAP = {
    "appear":              (1, 0, "set_only"),
    "fade-in":             (10, 0, "fade"),
    "wipe":                (22, 1, "wipe"),        # wipe(up):由下往上揭示
    "fly-in-from-bottom":  (2, 4, "fly"),
    "fly-in-from-left":    (2, 8, "fly"),
    "fly-in-from-right":   (2, 2, "fly"),
}
FLY_FROM = {  # (ppt_x 起點, ppt_y 起點)
    "fly-in-from-bottom": ("#ppt_x", "1+#ppt_h/2"),
    "fly-in-from-left":   ("0-#ppt_w/2", "#ppt_y"),
    "fly-in-from-right":  ("1+#ppt_w/2", "#ppt_y"),
}
TRANSITION_XML = {
    "fade": '<p:fade/>',
    "smooth-fade": '<p:fade/>',   # v1.0 相容
    "push": '<p:push dir="u"/>',
    "wipe": '<p:wipe dir="l"/>',
}
FX_WARNED = False


def hexc(s, default="0F172A"):
    return RGBColor.from_string((s or f"#{default}").lstrip("#"))


class DeckBuilder:
    def __init__(self, spec, json_dir):
        self.spec = spec
        self.json_dir = json_dir
        gs = spec["global_settings"]
        self.th = gs.get("theme_color", {})
        self.font_ea = gs.get("font_family_east_asian", "Microsoft JhengHei")
        self.font_latin = gs.get("font_family_latin", "Calibri")
        self.prs = Presentation()
        if gs.get("aspect_ratio", "16:9") == "4:3":
            self.prs.slide_width, self.prs.slide_height = Inches(10), Inches(7.5)
        else:
            self.prs.slide_width, self.prs.slide_height = Inches(13.333), Inches(7.5)
        self.pw = self.prs.slide_width / 914400  # 英吋
        self.ph = self.prs.slide_height / 914400

    # ---------- 低階工具 ----------

    def _style_run(self, run, size, bold, color):
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = self.font_latin
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", self.font_ea)

    def _textbox(self, slide, x, y, w, h, text, size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text or ""
        self._style_run(run, size, bold, color or hexc(self.th.get("primary")))
        return tb

    def _bg(self, slide, bg, default_hex):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                      self.prs.slide_width, self.prs.slide_height)
        rect.line.fill.background()
        rect.shadow.inherit = False
        fill = rect.fill
        if bg and bg.get("type") == "gradient":
            colors = bg.get("colors", [default_hex, default_hex])
            try:
                fill.gradient()
                stops = fill.gradient_stops
                stops[0].color.rgb = hexc(colors[0])
                stops[1].color.rgb = hexc(colors[-1])
                fill.gradient_angle = 45.0
            except Exception:
                fill.solid()
                fill.fore_color.rgb = hexc(colors[0])
        else:
            fill.solid()
            fill.fore_color.rgb = hexc((bg or {}).get("color", default_hex))
        return rect

    def _title_bar(self, slide, title):
        self._textbox(slide, 0.7, 0.35, self.pw - 1.4, 0.9, title,
                      size=28, bold=True)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.25),
            Inches(self.pw - 1.4), Emu(28575))  # 2.25pt 細線
        line.fill.solid()
        line.fill.fore_color.rgb = hexc(self.th.get("secondary"), "38BDF8")
        line.line.fill.background()
        line.shadow.inherit = False

    # ---------- 版型(回傳 [(shape, animation_dict), ...] 供動畫用) ----------

    def build_slide(self, sl, max_step=None):
        """max_step=None → 完整頁+動畫規格;max_step=k → 拆頁模式第 k 狀態(無動畫)。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        layout = sl.get("layout", "bullet_points")
        c = sl.get("content", {})
        anims = []

        def want(step):  # 拆頁過濾:無動畫元素永遠顯示
            return max_step is None or step is None or step <= max_step

        if layout in ("cover", "section_divider", "closing"):
            global FX_WARNED
            fx = sl.get("engine_config", {}).get("fx", "none")
            if fx not in ("none", None) and not FX_WARNED:
                print("[提示] native-pptx 引擎不支援 canvas 粒子特效(fx 已忽略);"
                      "需要粒子開場請用 build_html_deck.py 出封面")
                FX_WARNED = True
            self._bg(slide, c.get("background"),
                     self.th.get("primary", "#0F172A").lstrip("#"))
            white = RGBColor(0xFF, 0xFF, 0xFF)
            if layout == "section_divider":
                self._textbox(slide, 1, self.ph * 0.28, self.pw - 2, 1.2,
                              c.get("section_number", ""), size=54, bold=True,
                              color=hexc(self.th.get("accent"), "F59E0B"),
                              align=PP_ALIGN.CENTER)
                self._textbox(slide, 1, self.ph * 0.46, self.pw - 2, 1.5,
                              c.get("headline", ""), size=40, bold=True,
                              color=white, align=PP_ALIGN.CENTER)
            else:
                self._textbox(slide, 1, self.ph * 0.34, self.pw - 2, 1.6,
                              c.get("headline", ""), size=44, bold=True,
                              color=white, align=PP_ALIGN.CENTER)
                sub = c.get("sub_headline") or c.get("contact") or ""
                if sub:
                    self._textbox(slide, 1, self.ph * 0.56, self.pw - 2, 1.0,
                                  sub, size=20, bold=False,
                                  color=hexc(self.th.get("secondary"), "38BDF8"),
                                  align=PP_ALIGN.CENTER)

        elif layout == "bullet_points":
            self._title_bar(slide, sl.get("title", ""))
            y = 1.7
            for b in c.get("bullets", []):
                step = b.get("animation_step")
                if not want(step):
                    continue
                lv = b.get("level", 0)
                mark = "• " if lv == 0 else "– "
                tb = self._textbox(slide, 1.0 + lv * 0.6, y,
                                   self.pw - 2.0 - lv * 0.6, 0.75,
                                   mark + b.get("text", ""),
                                   size=20 if lv == 0 else 17,
                                   color=RGBColor(0x1F, 0x29, 0x37))
                if max_step is None and step is not None:
                    anims.append((tb, {"entry_order": step, "effect": "fade-in",
                                       "duration_ms": 400}))
                y += 0.78 if lv == 0 else 0.62

        elif layout == "two_column":
            self._title_bar(slide, sl.get("title", ""))
            colw = (self.pw - 2.0) / 2 - 0.3
            for i, side in enumerate(("left", "right")):
                col = c.get(side, {})
                x = 1.0 + i * (colw + 0.6)
                self._textbox(slide, x, 1.6, colw, 0.6, col.get("heading", ""),
                              size=20, bold=True,
                              color=hexc(self.th.get("secondary"), "38BDF8"))
                y = 2.35
                for b in col.get("bullets", []):
                    step = b.get("animation_step")
                    if not want(step):
                        continue
                    tb = self._textbox(slide, x, y, colw, 0.7,
                                       "• " + b.get("text", ""), size=17,
                                       color=RGBColor(0x1F, 0x29, 0x37))
                    if max_step is None and step is not None:
                        anims.append((tb, {"entry_order": step,
                                           "effect": "fade-in", "duration_ms": 400}))
                    y += 0.66

        elif layout == "architecture_chart":
            if c.get("background"):
                self._bg(slide, c.get("background"), "FFFFFF")
            self._title_bar(slide, sl.get("title", ""))
            nodes = {}
            for el in c.get("elements", []):
                anim = el.get("animation") or {}
                step = anim.get("entry_order")
                pos = el.get("position", {})
                if not want(step):
                    nodes[el.get("id")] = pos  # 連線仍需座標
                    continue
                st = el.get("style", {})
                if el.get("type", "shape_box") == "shape_box":
                    shp = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(float(pos.get("x", 0))), Inches(float(pos.get("y", 0))),
                        Inches(float(pos.get("w", 2))), Inches(float(pos.get("h", 1))))
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = hexc(st.get("fill", self.th.get("secondary", "#38BDF8")))
                    shp.line.fill.background()
                    tf = shp.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = el.get("text", "")
                    self._style_run(run, st.get("font_size", 18), False,
                                    hexc(st.get("color", "#FFFFFF")))
                else:
                    shp = self._textbox(slide, float(pos.get("x", 0)), float(pos.get("y", 0)),
                                        float(pos.get("w", 2)), float(pos.get("h", 0.5)),
                                        el.get("text", ""), size=st.get("font_size", 14),
                                        color=hexc(st.get("color", self.th.get("primary"))))
                nodes[el.get("id")] = pos
                if max_step is None and step is not None:
                    anims.append((shp, anim))
            for cn in c.get("connections", []):
                anim = cn.get("animation") or {}
                if anim.get("timeline_flow"):  # v1.0 相容
                    anim.setdefault("effect", "wipe")
                step = anim.get("entry_order")
                if not want(step):
                    continue
                a, b = nodes.get(cn.get("from")), nodes.get(cn.get("to"))
                if not a or not b:
                    continue
                ax = float(a.get("x", 0)) + float(a.get("w", 2))
                ay = float(a.get("y", 0)) + float(a.get("h", 1)) / 2
                bx = float(b.get("x", 0))
                by = float(b.get("y", 0)) + float(b.get("h", 1)) / 2
                if bx < ax:  # 反向:改走上下緣
                    ax = float(a.get("x", 0)) + float(a.get("w", 2)) / 2
                    ay = float(a.get("y", 0)) + float(a.get("h", 1))
                    bx = float(b.get("x", 0)) + float(b.get("w", 2)) / 2
                    by = float(b.get("y", 0))
                kind = MSO_CONNECTOR.ELBOW if cn.get("type") == "elbow" else MSO_CONNECTOR.STRAIGHT
                conn = slide.shapes.add_connector(kind, Inches(ax), Inches(ay),
                                                  Inches(bx), Inches(by))
                conn.line.color.rgb = hexc(self.th.get("secondary"), "38BDF8")
                conn.line.width = Pt(2.25)
                ln = conn.line._get_or_add_ln()
                tail = etree.SubElement(ln, qn("a:tailEnd"))
                tail.set("type", "triangle")
                if cn.get("label"):
                    self._textbox(slide, (ax + bx) / 2 - 0.8, (ay + by) / 2 - 0.55,
                                  1.6, 0.4, cn["label"], size=13, bold=True,
                                  color=hexc(self.th.get("primary")),
                                  align=PP_ALIGN.CENTER)
                if max_step is None and step is not None:
                    anims.append((conn, anim))

        elif layout == "figure_focus":
            self._title_bar(slide, sl.get("title", ""))
            p = self.json_dir / c.get("image_path", "")
            if p.is_file():
                slide.shapes.add_picture(str(p), Inches(1.2), Inches(1.7),
                                         height=Inches(self.ph - 2.8))
            else:
                self._textbox(slide, 1.2, 3.0, self.pw - 2.4, 1.0,
                              f"[圖檔未找到:{c.get('image_path')}]", size=18,
                              color=RGBColor(0xB9, 0x1C, 0x1C))
            if c.get("caption"):
                self._textbox(slide, 1.2, self.ph - 0.95, self.pw - 2.4, 0.6,
                              c["caption"], size=13,
                              color=RGBColor(0x4B, 0x55, 0x63), align=PP_ALIGN.CENTER)
        else:
            self._title_bar(slide, sl.get("title", f"(未知版型 {layout})"))

        if sl.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = sl["speaker_notes"]
        return slide, anims

    # ---------- 動畫 XML 注入 ----------

    def inject_timing(self, slide, anims):
        """anims: [(shape, anim_dict)] → <p:timing> 時間樹。entry_order 一組=一次點擊。"""
        if not anims:
            return
        groups = {}
        for shp, a in anims:
            if a.get("effect", "fade-in") not in EFFECT_MAP:
                print(f"[警告] 不支援的效果 {a.get('effect')},改用 fade-in")
                a = dict(a, effect="fade-in")
            groups.setdefault(int(a.get("entry_order", 1)), []).append((shp, a))
        self._ctn_id = 1

        def nid():
            self._ctn_id += 1
            return self._ctn_id

        def behaviors(spid, a):
            eff = a.get("effect", "fade-in")
            dur = int(a.get("duration_ms", 500))
            _, _, kind = EFFECT_MAP[eff]
            xml = (f'<p:set><p:cBhvr><p:cTn id="{nid()}" dur="1" fill="hold">'
                   f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
                   f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
                   f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                   f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>')
            if kind == "fade":
                xml += (f'<p:animEffect transition="in" filter="fade"><p:cBhvr>'
                        f'<p:cTn id="{nid()}" dur="{dur}"/>'
                        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>')
            elif kind == "wipe":
                xml += (f'<p:animEffect transition="in" filter="wipe(up)"><p:cBhvr>'
                        f'<p:cTn id="{nid()}" dur="{dur}"/>'
                        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>')
            elif kind == "fly":
                fx, fy = FLY_FROM[eff]
                for attr, frm in (("ppt_x", fx), ("ppt_y", fy)):
                    xml += (f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base">'
                            f'<p:cTn id="{nid()}" dur="{dur}" fill="hold"/>'
                            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
                            f'<p:attrNameLst><p:attrName>{attr}</p:attrName></p:attrNameLst></p:cBhvr>'
                            f'<p:tavLst><p:tav tm="0"><p:val><p:strVal val="{frm}"/></p:val></p:tav>'
                            f'<p:tav tm="100000"><p:val><p:strVal val="#{attr}"/></p:val></p:tav>'
                            f'</p:tavLst></p:anim>')
            return xml

        click_pars = []
        for order in sorted(groups):
            effect_pars = []
            for idx, (shp, a) in enumerate(groups[order]):
                spid = shp.shape_id
                pid, psub, _ = EFFECT_MAP[a.get("effect", "fade-in")]
                node_type = "clickEffect" if idx == 0 else "withEffect"
                delay = int(a.get("delay_ms", 0))
                effect_pars.append(
                    f'<p:par><p:cTn id="{nid()}" presetID="{pid}" presetClass="entr" '
                    f'presetSubtype="{psub}" fill="hold" grpId="0" nodeType="{node_type}">'
                    f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                    f'<p:childTnLst>{behaviors(spid, a)}</p:childTnLst></p:cTn></p:par>')
            click_pars.append(
                f'<p:par><p:cTn id="{nid()}" fill="hold">'
                f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>'
                f'<p:par><p:cTn id="{nid()}" fill="hold">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                f'<p:childTnLst>{"".join(effect_pars)}</p:childTnLst></p:cTn></p:par>'
                f'</p:childTnLst></p:cTn></p:par>')

        blds = "".join(f'<p:bldP spid="{shp.shape_id}" grpId="0"/>'
                       for shp, _ in anims)
        timing = (
            f'<p:timing xmlns:p="{NS_P}" xmlns:a="{NS_A}"><p:tnLst><p:par>'
            f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
            f'<p:cTn id="{nid()}" dur="indefinite" nodeType="mainSeq">'
            f'<p:childTnLst>{"".join(click_pars)}</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
            f'<p:bldLst>{blds}</p:bldLst></p:timing>')
        slide._element.append(etree.fromstring(timing))

    def inject_transition(self, slide, name):
        inner = TRANSITION_XML.get(name)
        if not inner:
            return
        xml = f'<p:transition xmlns:p="{NS_P}" spd="med">{inner}</p:transition>'
        el = etree.fromstring(xml)
        timing = slide._element.find(qn("p:timing"))
        if timing is not None:
            timing.addprevious(el)   # 規格順序:transition 在 timing 之前
        else:
            slide._element.append(el)


def collect_steps(sl):
    """回傳這頁所有 animation_step / entry_order 的排序清單。"""
    c = sl.get("content", {})
    steps = set()
    for b in c.get("bullets", []):
        if b.get("animation_step") is not None:
            steps.add(int(b["animation_step"]))
    for side in ("left", "right"):
        for b in c.get(side, {}).get("bullets", []) if isinstance(c.get(side), dict) else []:
            if b.get("animation_step") is not None:
                steps.add(int(b["animation_step"]))
    for el in c.get("elements", []):
        if (el.get("animation") or {}).get("entry_order") is not None:
            steps.add(int(el["animation"]["entry_order"]))
    for cn in c.get("connections", []):
        if (cn.get("animation") or {}).get("entry_order") is not None:
            steps.add(int(cn["animation"]["entry_order"]))
    return sorted(steps)


def main():
    ap = argparse.ArgumentParser(description="slides_content.json → 原生動畫 PPTX")
    ap.add_argument("json_file")
    ap.add_argument("-o", "--output", default="deck.pptx")
    ap.add_argument("--split", action="store_true",
                    help="動畫拆頁模式:每個動畫步驟拆成獨立頁面(轉 PDF 用)")
    args = ap.parse_args()

    json_path = Path(args.json_file)
    with open(json_path, encoding="utf-8") as f:
        spec = json.load(f)
    for key in ("metadata", "global_settings", "slides"):
        if key not in spec:
            sys.exit(f"[錯誤] slides_content.json 缺少頂層欄位 {key}")

    builder = DeckBuilder(spec, json_path.resolve().parent)
    skipped, total = [], 0
    for sl in sorted(spec["slides"], key=lambda s: s.get("slide_id", 0)):
        eng = sl.get("engine_config", {})
        if eng.get("engine") != "native-pptx":
            skipped.append(sl.get("slide_id"))
            continue
        split_this = args.split or eng.get("post_process") == "split-animation-layers"
        animate = eng.get("animate_objects", True)
        if split_this:
            steps = collect_steps(sl) if animate else []
            for k in [0] + steps:   # 狀態 0(未點擊)起,每步一頁
                builder.build_slide(sl, max_step=k)
                total += 1
        else:
            slide, anims = builder.build_slide(sl)
            if animate:
                builder.inject_timing(slide, anims)
            builder.inject_transition(slide, eng.get("transition"))
            total += 1
    if total == 0:
        sys.exit("[錯誤] 沒有任何頁面屬於 native-pptx 引擎;若整份走 html-canvas,請改用 build_html_deck.py")
    if skipped:
        print(f"[提示] 已跳過 html-canvas 頁面:slide_id {skipped}(用 build_html_deck.py 產出)")

    props = builder.prs.core_properties
    props.title = spec["metadata"].get("presentation_title", "")
    props.author = spec["metadata"].get("author", "")
    builder.prs.save(args.output)
    mode = "拆頁模式" if args.split else "原生動畫"
    print(f"[完成] {args.output}({total} 頁,{mode})")
    if not args.split:
        print("[提醒] 請用 PowerPoint 實際播放一次確認動畫;要轉 PDF 請加 --split 重出一份")


if __name__ == "__main__":
    main()
